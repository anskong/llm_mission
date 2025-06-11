### RAG Langchain 기반으로 산출물 문서를 읽어들인 후 Vector Embbeding의 Retreiver  LLM에  전달하여 
#   신규 문서 목차 별 내용을 LLM을 통해 생성하고  새로운 문서를 생성한다. 

# pip install langchain-core langchain-community langchain-text-splitters langchain-chroma langchain-openai  faiss-cpu
# pip install langchain-huggingface sentence-transformers
# pip install openai  # or use other embedding models
# pip install python-pptx openpyxl PyMuPDF
# pip install "unstructured[all-docs]"  # for PPT/Excel loader


# 1. 환경변수 읽기
import os
from dotenv import load_dotenv

import os
import json
import fitz  # PyMuPDF
import sys
from pathlib import Path
from tkinter import filedialog, Tk, simpledialog, messagebox
from pptx import Presentation
from pptx.util import Inches, Pt
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder, HumanMessagePromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.messages import SystemMessage,HumanMessage
from operator import itemgetter
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document
# from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_huggingface import HuggingFaceEmbeddings
# from langchain_huggingface.embeddings import HuggingFaceEmbeddings

from langchain_community.document_loaders import (
    PyPDFLoader,
    UnstructuredPowerPointLoader,
    UnstructuredExcelLoader,
    UnstructuredWordDocumentLoader
)

# from langchain_community.document_loaders import (
#     PyPDFLoader,
#     MSWordLoader,
#     MSPowerPointLoader,
#     UnstructuredExcelLoader,
# )

from langchain_openai import OpenAIEmbeddings

# 환경변수 읽어오기
load_dotenv(override=True)  # .env 파일을 덮어쓰기 모드로 읽기

# 환경변수 불러오기 (LLM API KEYS)
openai_key = os.getenv("OPENAI_API_KEY")
anthropic_key = os.getenv("ANTHROPIC_API_KEY")
huggingface_token = os.getenv("HUGGINGFACEHUB_API_TOKEN")

# ppt 생성 할 목차 및 목차 별 질문들
# contents_temp = [
#   {"title":"사업개요 - 추진배경","question":"""
#    바레인 사업의 추진배경은?
#    """},
#   {"title":"사업개요 - 추진목표","question":"바레인 사업의 추진목표는?"},
#   {"title":"제안요청사항 - 요구사항 총괄","question":"요구사항 총괄표를 작성해 주세요"},
#   {"title":"제안요청사항 - 요구사항 목록","question":"요구사항 목록표를 작성해 주세요"},
#   # {"title":"사업개요 - 추진배경","question":"바레인 사업의 추진배경은?"},
# ]

# embedding_model_name = "sentence-transformers/all-MiniLM-L6-v2"
# embedding_model_name = "jhgan/ko-sroberta-multitask"
embedding_model_name = r"C:\ai_dev\ko-sroberta-multitask"

# [사용할 LLM 인스턴스 생성]

llm = ChatOpenAI(
  openai_api_base="http://localhost:1234/v1",
  openai_api_key="lm-studio",
  model_name="exaone-3.5-2.4b-instruct",
  temperature=0.7,
)

# 최종 결과내용을 저장할 ppt 명칭 및 경로
output_path = None

def main():
  # response = llm.invoke("너는 누구니?")
  # print("llm connection test ::: ",response.content)
  
  # 파일 선택기 열기
  root = Tk()
  root.withdraw()

  contents_path = filedialog.askopenfilename(title="TXT 파일 선택", filetypes=[("TXT files", "*.txt")])

  if not contents_path:
    messagebox.showwarning("경고", "목차 파일을 선택하지 않았습니다.")
    return
  else :
    print(f"선택한 파일 명 :::: {contents_path}")
    # contents_path = "contents.txt"
    with open(contents_path, 'r', encoding='utf-8') as file:
        global contents
        contents = json.load(file)
        print(f"목차 내용 ::: {contents}")


  pdf_path = filedialog.askopenfilename(title="PDF 파일 선택", filetypes=[("PDF files", "*.pdf")])

  if not pdf_path:
    messagebox.showwarning("경고", "PDF 파일을 선택하지 않았습니다.")
    return
  else :
    print(f"선택한 파일 명 :::: {pdf_path}")
  

  output_path = filedialog.asksaveasfilename(defaultextension=".pptx",
                                               filetypes=[("PowerPoint files", "*.pptx")],
                                               title="저장할 PPT 파일 이름")
  if not output_path:
        messagebox.showwarning("경고", "저장 파일명을 입력하지 않았습니다.")
        return
  else :
      print(f"저장 할 파일 명 first :::: {output_path}")


  # PDF 문서 로딩 및 chunk 분할
  loader = PyPDFLoader(pdf_path)
  pages = loader.load()
  print(f"문서 로딩한 pages 수 :::: {len(pages)}")

  splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=300)
  splits = splitter.split_documents(pages)
  print(f"문서를 chunk 분할한 수 :::: {len(splits)}")

  # Vector Embedding 및 Retreiver 인스턴스 생성
  # embedding = OpenAIEmbeddings(model="text-embedding-3-small", openai_api_key=openai_key)


  #HuggingfaceEmbedding 함수로 Open source 임베딩 모델 로드
#   model_name = "jhgan/ko-sroberta-multitask"
  
  ko_embedding= HuggingFaceEmbeddings(
      model_name=embedding_model_name
  )

  vectorstore = FAISS.from_documents(splits, ko_embedding)
  # vector store 검색시 유사문서는 5개만 반환하라
  retriever = vectorstore.as_retriever(search_kwargs={"k": 5})

  # 프롬프트 템플릿
  prompt_template = ChatPromptTemplate.from_messages([
    SystemMessage("""
                  당신은 건강보험 업무 전문가이며 문서 작성에도 탁월합니다.
                  제공되는 컨텍스트를 최대한 활용하여 바레인 제안서 ppt를 세부적으로 작성해 주세요
                  """),
    MessagesPlaceholder("chat_history"),
    HumanMessagePromptTemplate.from_template(
        """
        주어진 컨텍스트에 따라 다음 질문에 답하십시오.\n컨텍스트：{context}\n질문：{question}
        질문에 대한 답변은 마크다운 형식으로 자세하게 정리하고 질문에 표함 된 예시는 답변 형식에 참조하세요
        """
    )
  ])
  # 출력 파서 정의
  parser = StrOutputParser()
  
  # 11) LCEL 기반 RAG 체인 구성
  rag_chain = (
    {
        "context": itemgetter("question") | retriever | format_docs,
        "question": itemgetter("question"),
        "chat_history": itemgetter("chat_history"),
    }
    | prompt_template
    | llm
    | parser
  )

  # 12) 체인 실행
  history = []
  # question = "바레인 사업 추진 목표는?"
  # answer = rag_chain.invoke({"question": question, "chat_history": history})

  # 13) 체인 실행 및 문서 생성
  generate_doc_from_llm(rag_chain,history,output_path)

  # 13) 결과 출력
  # print(f"\n🧠 질문: {question}")
  # print("📝 답변:\n", answer)


def format_docs(docs):
    return "\n\n".join(doc.page_content for doc in docs)

# contents에 따라 RAG 체인을 실행하고 llm 결과를 slides[]에 저장하여 ppt 생성에 전달
def generate_doc_from_llm(rag_chain,history,output_path):
  slides = []
  for content in contents:
    question = content['question']
    print(f"llm question ::: {question}")
    answer = rag_chain.invoke({"question": question, "chat_history": history})
    slides.append((content['title'],answer))
  
  print("🎞 PPT 생성 중...")
  create_ppt(slides,output_path)
  print("✅ 완료! 저장 위치:", output_path)

def create_ppt(slide_contents,output_path):
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "목차"
    content = title_slide.placeholders[1].text_frame

    for idx, (title, _) in enumerate(slide_contents, 1):
        content.add_paragraph().text = f"{idx}. {title}"

    for index, (title, content_text) in enumerate(slide_contents):
        chunks = split_text_to_slides(content_text)
        for i, chunk in enumerate(chunks):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide_title = title if i == 0 else f"{title} (계속)"
            slide.shapes.title.text = slide_title

            para = slide.shapes.title.text_frame.paragraphs[0]
            if para.runs:
                run = para.runs[0]
                run.font.size = Pt(18)
                run.font.name = "맑은 고딕"

            content_box = slide.placeholders[1]
            content_box.text = chunk
            for p in content_box.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.name = "맑은 고딕"

            # matched_image = match_image_by_index(index, page_images)
            # if matched_image:
            #     slide.shapes.add_picture(matched_image, Inches(5.5), Inches(1.5), width=Inches(3))
    print(f"저장 할 파일 명 second :::: {output_path}")
    prs.save(output_path)    
    
def split_text_to_slides(text, max_chars=800):
    paragraphs = text.split('\n')
    slides = []
    current_slide = ""

    for para in paragraphs:
        if len(current_slide) + len(para) + 1 < max_chars:
            current_slide += para + '\n'
        else:
            slides.append(current_slide.strip())
            current_slide = para + '\n'
    if current_slide:
        slides.append(current_slide.strip())

    return slides

# ✅ GUI로 여러 파일 선택
def select_multiple_files():
    root = Tk()
    root.withdraw()
    file_paths = filedialog.askopenfilenames(
        title="PDF, PPTX, XLSX,DOCX 파일 선택",
        filetypes=[
            ("Supported files", "*.pdf *.pptx *.xlsx *.docx"),
            ("All files", "*.*"),
        ],
    )
    return list(file_paths)

# ✅ 확장자에 따른 LangChain Loader 선택
def load_documents(file_paths):
    all_docs = []
    for path in file_paths:
        ext = os.path.splitext(path)[1].lower()
        if ext == ".pdf":
            loader = PyPDFLoader(path)
        elif ext == ".pptx":
            loader = UnstructuredPowerPointLoader(path)
            # loader = MSWordLoader(path)
        elif ext == ".xlsx":
            loader = UnstructuredExcelLoader(path)
        elif ext == ".docx":
            loader = UnstructuredWordDocumentLoader(path)
        else:
            print(f"⚠️ 지원하지 않는 형식: {ext}")
            continue
        docs = loader.load()
        all_docs.extend(docs)
    return all_docs

# ✅ 문서 분할
def split_documents(documents,chunk_size=1000, chunk_overlap=200):
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    return splitter.split_documents(documents)

# ✅ 벡터스토어에 저장 (FAISS)
def embed_and_store(documents, persist_path="faiss_index"):
    # embeddings = OpenAIEmbeddings()  # OpenAI API 키 필요 (환경변수 OPENAI_API_KEY)
    
    # model_name = "jhgan/ko-sroberta-multitask"

    print(f"embedding_model_name ::: {embedding_model_name}")

    ko_embedding= HuggingFaceEmbeddings(
        model_name=embedding_model_name
    )

    vectorstore = FAISS.from_documents(documents, ko_embedding)
    vectorstore.save_local(persist_path)
    print(f"✅ FAISS 벡터스토어 저장 완료: {persist_path}")
    return vectorstore


# ✅ (사용자가 작성한) 목차/질문 파일 열고 읽어들이기
def open_contents_file():
    root = Tk()
    root.withdraw()
    contents_path = filedialog.askopenfilename(title="TXT 파일 선택", filetypes=[("TXT files", "*.txt")])
    
    if not contents_path:
        messagebox.showwarning("경고", "목차 파일을 선택하지 않았습니다.")
        return None
    else:
        print(f"선택한 파일 명 :::: {contents_path}")
        with open(contents_path, 'r', encoding='utf-8') as file:
            global contents
            contents = json.load(file)
            # print(f"목차 내용 ::: {contents}")
            return contents

def select_output_path():
    root = Tk()
    root.withdraw()
    output_path_ = filedialog.asksaveasfilename(defaultextension=".pptx",
                                               filetypes=[("PowerPoint files", "*.pptx")],
                                               title="저장할 PPT 파일 이름")
    if not output_path_:
        messagebox.showwarning("경고", "저장 파일명을 입력하지 않았습니다.")
        exit()
    else:
        print(f"저장 할 파일 명 :::: {output_path_}")
        return output_path_

# ✅ 전체 워크플로우
def main_workflow():
    
    print("📂 목차 파일을 선택하세요...")
    open_contents_file()
    print(f"목차 내용 ::: {contents}")

    print("📂 문서 파일을 선택하세요...")
    file_paths = select_multiple_files()
    if not file_paths:
        print("❗파일이 선택되지 않았습니다.")
        exit()
    print(f"📂 선택된 파일: {file_paths}")

    print("📂 결과를 저장할 PPT 파일 경로를 선택하세요...")
    output_path = select_output_path()

    documents = load_documents(file_paths)
    print(f"📄 총 로딩된 문서 수: {len(documents)}")

    split_docs = split_documents(documents, chunk_size=1000, chunk_overlap=200)
    if not split_docs:
        print("❗문서 분할에 실패했습니다.")
        exit()
    print(f"✂️ 분할된 청크 수: {len(split_docs)}")

    vectorstore = embed_and_store(split_docs)
    retriever = vectorstore.as_retriever(search_kwargs={"k": 5})

    # 프롬프트 템플릿
    prompt_template = ChatPromptTemplate.from_messages([
        SystemMessage("""
                    당신은 건강보험 업무 전문가이며 문서 작성에도 탁월합니다.
                    제공되는 컨텍스트를 최대한 활용하여 바레인 제안서 ppt를 세부적으로 작성해 주세요
                    """),
        MessagesPlaceholder("chat_history"),
        HumanMessagePromptTemplate.from_template(
            """
            주어진 컨텍스트에 따라 다음 질문에 답하십시오.\n컨텍스트：{context}\n질문：{question}
            질문에 대한 답변은 마크다운 형식으로 자세하게 정리하고 질문에 표함 된 예시는 답변 형식에 참조하세요
            """
        )
    ])
    # 출력 파서 정의
    parser = StrOutputParser()
    
    # LCEL 기반 RAG 체인 구성
    rag_chain = (
        {
            "context": itemgetter("question") | retriever | format_docs,
            "question": itemgetter("question"),
            "chat_history": itemgetter("chat_history"),
        }
        | prompt_template
        | llm
        | parser
    )

    # 12) 체인 실행
    history = []
    # question = "바레인 사업 추진 목표는?"
    # answer = rag_chain.invoke({"question": question, "chat_history": history})

    # 13) 체인 실행 및 문서 생성
    generate_doc_from_llm(rag_chain,history,output_path)



if __name__ == "__main__":
    main_workflow()
    
  
# if __name__ == "__main__":
#   main()

